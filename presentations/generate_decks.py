#!/usr/bin/env python3
"""Generate KlimaGuard Kids pitch and check-in PowerPoint decks."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent

# Brand palette (matches src/app/globals.css)
OCEAN = RGBColor(0x0E, 0xA5, 0xE9)
INK = RGBColor(0x0F, 0x17, 0x2A)
SKY = RGBColor(0xE0, 0xF2, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x47, 0x55, 0x69)
LEAF = RGBColor(0x4A, 0xDE, 0x80)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, OCEAN)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(8.8), Inches(1.5))
    stf = sub.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(20)
    sp.font.color.rgb = SKY
    sp.alignment = PP_ALIGN.CENTER


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(8.8), Inches(1.0))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.9), Inches(0.7))
    ttf = title_box.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = OCEAN

    top = 1.25 if subtitle else 1.15
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(8.9), Inches(0.5))
        stf = sub_box.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED
        top = 1.55

    body = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(8.7), Inches(5.0))
    btf = body.text_frame
    btf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = INK
        p.space_after = Pt(10)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.9), Inches(0.7))
    ttf = title_box.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = OCEAN

    for col, (ctitle, bullets, x) in enumerate(
        [(left_title, left_bullets, 0.55), (right_title, right_bullets, 5.05)]
    ):
        h = slide.shapes.add_textbox(Inches(x), Inches(1.15), Inches(4.2), Inches(0.4))
        htf = h.text_frame
        htf.text = ctitle
        hp = htf.paragraphs[0]
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = INK

        body = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(4.2), Inches(4.8))
        btf = body.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = MUTED
            p.space_after = Pt(8)


def add_metrics_slide(prs: Presentation, title: str, metrics: list[tuple[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SKY)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.9), Inches(0.7))
    ttf = title_box.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = INK

    cols = min(len(metrics), 4)
    width = 8.4 / cols
    for i, (value, label) in enumerate(metrics):
        x = 0.55 + i * width
        card = slide.shapes.add_shape(
            1, Inches(x), Inches(1.8), Inches(width - 0.15), Inches(3.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = OCEAN

        vbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.3), Inches(width - 0.35), Inches(1.2))
        vtf = vbox.text_frame
        vtf.text = value
        vp = vtf.paragraphs[0]
        vp.font.size = Pt(32)
        vp.font.bold = True
        vp.font.color.rgb = OCEAN
        vp.alignment = PP_ALIGN.CENTER

        lbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.5), Inches(width - 0.35), Inches(1.5))
        ltf = lbox.text_frame
        ltf.text = label
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.font.size = Pt(12)
        lp.font.color.rgb = MUTED
        lp.alignment = PP_ALIGN.CENTER


def build_pitch_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "KlimaGuard Kids",
        "Child-centric climate health intelligence · Open Source · Agentic AI",
    )

    add_bullet_slide(
        prs,
        "The challenge",
        [
            "Climate change is reshaping children's health today — not in the distant future.",
            "Heatwaves, floods, air pollution, vectors, and food insecurity compound fastest in LMICs.",
            "2.4B children need anticipatory, age-appropriate guidance families can act on within days.",
            "Most climate tools speak to policymakers — not to children, caregivers, or schools.",
        ],
    )

    add_bullet_slide(
        prs,
        "Our solution",
        [
            "Open-source agentic AI platform connecting climate data with health intelligence.",
            "Nine cooperating agents ingest trusted feeds and output unified preparedness briefings.",
            "Age-banded guidance for children 5–17, plus caregiver and teacher action steps.",
            "Privacy-first: country/city location only; no child accounts required for demo.",
        ],
    )

    add_section_slide(prs, "How it works")

    add_bullet_slide(
        prs,
        "Nine AI agents — one pipeline",
        [
            "Climate Data Agent — live 7-day forecast, heat index, AQI (Open-Meteo)",
            "Health Risk Agent — heat, respiratory, flood, vector stressors (WHO-aligned)",
            "Nutrition Agent — hydration, food safety, scarcity patterns",
            "Disease Agent — transmission pathways, precautions, illness profiles",
            "Natural Medicine Agent — evidence-tagged remedies under caregiver supervision",
            "Synthesis Agent — cross-correlation + age-banded child guidance",
            "India Regional + Impact Agents — CHIS scores across 12 Indian regions",
            "Kids Health Chat Agent — age-banded Q&A with privacy safeguards",
        ],
    )

    add_two_column_slide(
        prs,
        "Technology stack",
        "Platform",
        [
            "Next.js 15 · React 19 · TypeScript",
            "Tailwind CSS 4 · REST API",
            "Live Open-Meteo weather & air quality",
            "MIT open-source license",
        ],
        "Delivery",
        [
            "Web dashboard + India impact dashboard + kids health chat",
            "POST /api/analyze — full agent pipeline",
            "POST /api/chat — age-banded health Q&A",
            "Roadmap: PWA, SMS/USSD, 12+ languages",
        ],
    )

    add_metrics_slide(
        prs,
        "Prototype traction",
        [
            ("20", "Countries in demo registry"),
            ("9", "Cooperating AI agents"),
            ("12", "India regions with CHIS"),
            ("<5s", "End-to-end analysis per location"),
        ],
    )

    add_bullet_slide(
        prs,
        "Who we serve",
        [
            "Children, caregivers, teachers, and community health workers",
            "Local health and education authorities",
            "Global coverage — extensible to any coordinates",
            "Built for local adaptation under MIT license",
        ],
    )

    add_bullet_slide(
        prs,
        "Impact & SDG alignment",
        [
            "SDG 3 — Anticipatory action for climate-sensitive health risks",
            "SDG 13 — Child-centred climate early warning",
            "SDG 2 — Nutrition & food security inference",
            "SDG 4 — School-ready heat and flood guidance",
        ],
    )

    add_bullet_slide(
        prs,
        "Safeguarding & data ethics",
        [
            "No child accounts required for demo",
            "Country/city-level location only",
            "Agents cite data provenance on every report",
            "Roadmap: COPPA/GDPR-K, in-region language models, national met feeds",
        ],
    )

    add_section_slide(prs, "12-month roadmap")

    add_bullet_slide(
        prs,
        "Milestones (next 12 months)",
        [
            "Q1: Field-pilot design · 3 languages · usability testing",
            "Q2: Structured pilot (500+ users) · offline PWA · WHO/met integrations",
            "Q3: SMS/USSD · 8 more languages · privacy audit · open-source docs",
            "Q4: Pilot evaluation · ministry partnerships · multi-country rollout plan",
        ],
    )

    add_title_slide(
        prs,
        "Try the live demo",
        "Dashboard → Select country → Run agent analysis → Child guidance",
    )

    return prs


def build_checkin_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "KlimaGuard Kids — Check-In",
        "Progress update · May 2026",
    )

    add_bullet_slide(
        prs,
        "Check-in agenda",
        [
            "Progress since last review",
            "Technical status & demo",
            "Prototyping & testing results",
            "Blockers & decisions needed",
            "Next 90 days",
        ],
    )

    add_metrics_slide(
        prs,
        "Status at a glance",
        [
            ("Live", "Deployed prototype"),
            ("20", "Countries supported"),
            ("9", "Agents operational"),
            ("100%", "Build passing"),
        ],
    )

    add_two_column_slide(
        prs,
        "Completed this period",
        "Product",
        [
            "Home, dashboard, India, chat, and about pages",
            "Global country selector (20 countries)",
            "Live Open-Meteo integration",
            "Age-banded child guidance output",
        ],
        "Engineering",
        [
            "Nine-agent orchestration pipeline",
            "REST API (/api/analyze, /api/chat, /api/countries)",
            "Agent provenance & status UI",
            "Production build verified",
        ],
    )

    add_bullet_slide(
        prs,
        "Prototyping & testing results",
        [
            "Quantitative: End-to-end runs complete in seconds; pipeline stable across diverse climates.",
            "Qualitative: Unified briefings rated clearer than siloed weather/health alerts in demo reviews.",
            "Gaps: No formal field pilot yet; comprehension and behaviour-change metrics pending.",
            "Next: Structured pilot with schools and community health workers.",
        ],
    )

    add_two_column_slide(
        prs,
        "Risks & blockers",
        "Current blockers",
        [
            "Field partners not yet signed",
            "Localization not started",
            "Offline / SMS channels not built",
            "Formal privacy audit pending",
        ],
        "Mitigations",
        [
            "Outreach to MoH / education ministries",
            "Prioritize 3 pilot languages in Q1",
            "PWA scoped for Q2 delivery",
            "COPPA/GDPR-K review in Q3",
        ],
    )

    add_bullet_slide(
        prs,
        "Next 90 days",
        [
            "Finalize 2 field-pilot partner sites",
            "Run usability tests with caregivers and teachers (n≈30)",
            "Ship 3 priority languages",
            "Integrate first national met-service or WHO outbreak feed",
            "Publish pilot evaluation plan and success metrics",
        ],
    )

    add_bullet_slide(
        prs,
        "Decisions requested",
        [
            "Confirm pilot geography and partner institutions",
            "Approve success metrics (comprehension, time-to-action, confidence)",
            "Align on data-sharing and safeguarding requirements",
            "Confirm resourcing for localization and offline PWA (Q2)",
        ],
    )

    add_title_slide(
        prs,
        "Questions & next steps",
        "Demo available · Open to feedback · Thank you",
    )

    return prs


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    pitch_path = OUT_DIR / "KlimaGuard-Kids-Pitch.pptx"
    checkin_path = OUT_DIR / "KlimaGuard-Kids-Check-In.pptx"

    build_pitch_deck().save(pitch_path)
    build_checkin_deck().save(checkin_path)

    print(f"Created: {pitch_path}")
    print(f"Created: {checkin_path}")


if __name__ == "__main__":
    main()
